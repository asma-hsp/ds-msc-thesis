"""Build the defense deck on the UniPD template.
Output: slides/defense_slides.pptx. Run make_charts.py and make_slide_assets.py first.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = "/extra/ahoseinp/projects/ds-msc-thesis"
FIG = f"{REPO}/figures"
ASSETS = f"{REPO}/slides/assets"
OUT = f"{REPO}/slides/defense_slides.pptx"

INK = RGBColor(0x0B, 0x0B, 0x0B)
MUTED = RGBColor(0x52, 0x51, 0x4E)
RED = RGBColor(0x9B, 0x00, 0x13)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
LIGHT = RGBColor(0xF0, 0xEF, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROSE = RGBColor(0xFB, 0xE9, 0xE7)

IMG = {
    "tintin_fig": f"{ASSETS}/icir_tintin.png",
    "inst_cat": f"{ASSETS}/instance_vs_category.png",
    "clip": f"{ASSETS}/clip_schematic.png",
    "stats": f"{ASSETS}/icir_stats_trim.png",
    "basic": f"{ASSETS}/basic_method_trim.png",
    "harris": f"{FIG}/harris_surface.png",
    "success": f"{ASSETS}/qual_success_trim.png",
    "failure": f"{ASSETS}/qual_failure_trim.png",
    "ladder_l": f"{ASSETS}/ladder_clipl.png",
    "ladder_s2": f"{ASSETS}/ladder_siglip2.png",
    "backbones": f"{ASSETS}/backbones.png",
    "ref": f"{ASSETS}/chair_ref_43.jpg",
    "top1": f"{ASSETS}/chair_set_43.jpg",
    "paper_stats": f"{ASSETS}/icir_paper_stats.png",
    "pipeline": f"{ASSETS}/pipeline_thesis.png",   # compiled from figures/pipeline_overview.tex
    "clip_siglip": f"{ASSETS}/clip_vs_siglip.png",
    "qual_success": f"{ASSETS}/qual_homer_success_q832.png",
    "qual_failure": f"{ASSETS}/qual_dogglasses_failure_q1490.png",
    "tintin_q": f"{FIG}/caption_example_samples/ref.jpg",
    "eq_pk": f"{ASSETS}/eq_pk.png",
    "eq_ap": f"{ASSETS}/eq_ap.png",
    "eq_map": f"{ASSETS}/eq_map.png",
    "eq_mmap": f"{ASSETS}/eq_mmap.png",
    "eq_centering": f"{ASSETS}/eq_centering.png",
    "eq_minnorm": f"{ASSETS}/eq_minnorm.png",
    "eq_harris": f"{ASSETS}/eq_harris.png",
    "eq_triplet": f"{ASSETS}/eq_triplet.png",
    "eq_problem": f"{ASSETS}/eq_problem.png",
}

prs = Presentation(f"{REPO}/slides/UniPD_template.pptx")
CONTENT_LAYOUT = prs.slide_masters[0].slide_layouts[1]

# ---------------- helpers ----------------
def style_runs(paragraph, size, bold=False, color=INK, italic=False):
    for r in paragraph.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color

def add_pic(slide, img, left, top, max_w, max_h, border=None):
    w, h = Image.open(img).size
    scale = min(max_w / w, max_h / h)
    dw, dh = int(w * scale), int(h * scale)
    pic = slide.shapes.add_picture(img, int(left + (max_w - dw) / 2),
                                   int(top + (max_h - dh) / 2), dw, dh)
    if border is not None:
        pic.line.color.rgb = border
        pic.line.width = Pt(1.5)
    return pic

def txt(slide, left, top, w, h, text, size=17, bold=False, color=INK,
        align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        style_runs(p, size, bold, color, italic)
    return box

def clean_placeholders(slide):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx != 0:
            ph._element.getparent().remove(ph._element)

def new_slide(title, notes=""):
    slide = prs.slides.add_slide(CONTENT_LAYOUT)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        style_runs(p, 25, bold=True, color=WHITE)   # red banner -> white title
    clean_placeholders(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    n = len(prs.slides._sldIdLst)
    txt(slide, Inches(8.9), Inches(6.94), Inches(0.75), Inches(0.32), str(n),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)   # white on footer band
    return slide

def bullets(slide, items, left=Inches(0.68), top=Inches(1.55), w=Inches(8.6),
            size=17, gap=8):
    h = Inches(6.75) - top
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, opts = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if opts.get("indent"):
            p.text = "      –  " + text
        elif opts.get("nobullet"):
            p.text = text
        else:
            p.text = "•  " + text
        p.space_after = Pt(gap)
        style_runs(p, opts.get("size", size), opts.get("bold", False),
                   opts.get("color", INK), opts.get("italic", False))
    return box

def takeaway(slide, text, top):
    lines = text.split("\n")
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.68), top,
                                 Inches(8.6), Inches(0.55 if len(lines) == 1 else 0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = LIGHT
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = PP_ALIGN.CENTER
        style_runs(p, 15, bold=True, color=RED)

def card(slide, left, top, w, h, text=None, fill=WHITE, line=MUTED, size=13,
         bold=False, color=INK, italic=False, align=PP_ALIGN.CENTER):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    if line is None:
        c.line.fill.background()
    else:
        c.line.color.rgb = line
        c.line.width = Pt(1)
    if text is not None:
        tf = c.text_frame
        tf.word_wrap = True
        for i, ln in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ln
            p.alignment = align
            style_runs(p, size, bold, color, italic)
    return c

def arrow(slide, left, top, w, h=Inches(0.42)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = RED
    a.line.fill.background()
    return a

def table(slide, rows, left, top, width, row_h=0.42, size=14,
          highlight_rows=(), col0_left=True):
    n, m = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(n, m, left, top, width,
                                 Inches(row_h * n)).table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = (PP_ALIGN.LEFT if (j == 0 and col0_left)
                               else PP_ALIGN.CENTER)
                style_runs(p, size, bold=(i == 0 or i in highlight_rows),
                           color=WHITE if i == 0 else INK)
    return tbl

# =========================================================
# 1 — title
# =========================================================
s = prs.slides[0]   # solid UniPD-red background -> everything in white
tp = s.shapes.title
tp.left, tp.top, tp.width, tp.height = (Inches(0.55), Inches(1.0),
                                        Inches(8.9), Inches(2.15))
tf = tp.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.paragraphs[0].text = "Imagining the Target:"
p = tf.add_paragraph()
p.text = ("Zero-Shot Instance-Level Composed Image Retrieval "
          "with MLLM-Generated Captions")
p.space_before = Pt(10)
for i, par in enumerate(tf.paragraphs):
    par.alignment = PP_ALIGN.CENTER
    style_runs(par, 30 if i == 0 else 21, bold=(i == 0), color=WHITE)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.25), Inches(3.38),
                          Inches(1.5), Inches(0.03))
rule.fill.solid(); rule.fill.fore_color.rgb = WHITE
rule.line.fill.background()
for ph in s.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.left, ph.top, ph.width, ph.height = (Inches(1.25), Inches(3.65),
                                                Inches(7.5), Inches(2.2))
        tf2 = ph.text_frame
        tf2.word_wrap = True
        sub = [("Asma Hoseinpour Siouki", 17, True),
               ("MSc in Data Science — University of Padova", 13, False),
               ("", 6, False),
               ("Supervisor: Prof. Lamberto Ballan", 12.5, False),
               ("Co-supervisor: Prof. Marco Fiorucci", 12.5, False),
               ("July 2026", 12, False)]
        for i, (t, sz, bd) in enumerate(sub):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = t
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(4)
            style_runs(p, sz, bold=bd, color=WHITE)
s.notes_slide.notes_text_frame.text = (
    "[0:20]  Good morning everyone, and thank you for being here.")
sldIdLst = prs.slides._sldIdLst
prs.part.drop_rel(sldIdLst[1].rId)
sldIdLst.remove(sldIdLst[1])

# =========================================================
# 2 — Composed Image Retrieval
# =========================================================
s = new_slide("Composed Image Retrieval", notes=(
    "[1:05]  Imagine that you see a chair online that you really like, but you would "
    "like to find a dining set that includes that same chair design together with a "
    "table.\n\n"
    "Searching only with the IMAGE of the chair would mainly return the same chair or "
    "other visually similar chairs. It would not understand that you are looking for it "
    "as part of a dining set.\n\n"
    "On the other hand, searching only for 'a dining table and chairs' would return "
    "many different sets, but probably not the specific chair you selected.\n\n"
    "To express exactly what you are looking for, you need to combine both elements: "
    "the image identifies the chair, while the text describes the context in which you "
    "want to find it.\n\n"
    "This is an example of composed image retrieval, the task I worked on in my thesis. "
    "The query consists of a reference image and a modification text. The goal is to "
    "retrieve images that preserve the information in the reference image while also "
    "satisfying the modification described by the text.\n\n"
    "One naming note: the reference image and modification text are the two parts of the "
    "query, so from here on I will simply call them the QUERY IMAGE and the QUERY TEXT."))
Y, BH = Inches(1.7), Inches(2.0)
add_pic(s, IMG["ref"], Inches(0.68), Y, Inches(2.5), BH, border=MUTED)
txt(s, Inches(3.2), Y, Inches(0.4), BH, "+", size=30, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
card(s, Inches(3.6), Inches(2.25), Inches(2.0), Inches(0.9),
     text="“placed around\na table”", size=13, italic=True, color=BLUE)
arrow(s, Inches(5.7), Inches(2.5), Inches(0.6))
add_pic(s, IMG["top1"], Inches(6.4), Y, Inches(2.5), BH, border=MUTED)
txt(s, Inches(0.68), Inches(3.75), Inches(2.5), Inches(0.4),
    "reference image\n(query image)", size=12, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, Inches(3.6), Inches(3.25), Inches(2.0), Inches(0.4),
    "modification text\n(query text)", size=12, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, Inches(6.4), Inches(3.75), Inches(2.5), Inches(0.35), "target image",
    size=12, color=MUTED, align=PP_ALIGN.CENTER)
bullets(s, [
    ("Image only  →  the same chair, or visually similar chairs — but not as a dining set", {}),
    ("Text only  →  many dining sets — but not your specific chair", {}),
    ("Both together: the image identifies the object, the text describes the context "
     "in which you want to find it", {"bold": True}),
], top=Inches(4.7), size=16, gap=7)
txt(s, Inches(0.68), Inches(4.22), Inches(8.6), Inches(0.32),
    "From here on we call them the query image and the query text.",
    size=11.5, italic=True, color=RED, align=PP_ALIGN.CENTER)
takeaway(s, "Goal: retrieve images that preserve the query object and satisfy the query text.",
         Inches(6.15))

# =========================================================
# 3 — instance-level
# =========================================================
s = new_slide("Instance-level composed retrieval", notes=(
    "[0:55]  My thesis focuses specifically on the INSTANCE-LEVEL version of this task.\n\n"
    "In category-level composed image retrieval, any object belonging to the correct "
    "class may satisfy the query. For example, if the query asks for a chair around a "
    "dining table, any suitable chair design might be considered correct.\n\n"
    "In instance-level retrieval, this is not enough. The system must retrieve the exact "
    "same chair shown in the query image, while also satisfying the query text.\n\n"
    "Therefore, the main challenge is to DETECT that specific instance among many "
    "look-alikes of the same class — even when its context, appearance, viewpoint, or "
    "visual domain changes."))
add_pic(s, IMG["inst_cat"], Inches(0.55), Inches(1.45), Inches(8.9), Inches(3.9))
bullets(s, [
    ("Category-level:  any object of the correct class may satisfy the query", {}),
    ("Instance-level:  the system must retrieve the exact same object shown in the "
     "query image, while also satisfying the query text", {"bold": True}),
    ("The challenge: detect that one instance among many look-alikes of the same class — "
     "even when its context, appearance, viewpoint, or visual domain changes", {}),
], top=Inches(5.4), size=15, gap=5)

# =========================================================
# 4 — i-CIR benchmark (+ scale strip)
# =========================================================
s = new_slide("The i-CIR benchmark", notes=(
    "[1:05]  For our experiments we use the i-CIR benchmark, specifically designed for "
    "instance-level composed image retrieval.\n\n"
    "Each object instance is associated with several composed queries, and every "
    "composed query consists of a query image and a query text.\n\n"
    "Importantly, each query may have MULTIPLE correct target images — several database "
    "images can correctly depict the same object under the requested modification. So "
    "evaluation is not about finding one right answer, but ranking all of them highly.\n\n"
    "In addition, each instance has its own database of carefully curated hard "
    "negatives. These are designed to confuse the retrieval system: some are visually "
    "similar to the query image but fail to satisfy the query text; others satisfy "
    "the query text but depict a different object instance; some match both parts "
    "partially while still being incorrect.\n\n"
    "This prevents the system from succeeding by relying only on the image or only on "
    "the text.\n\n"
    "[Scale and full distributions are on the next slide.]"))
add_pic(s, IMG["tintin_fig"], Inches(0.55), Inches(1.4), Inches(8.9), Inches(3.95))
bullets(s, [
    ("Each instance → several composed queries (query image + query text)", {}),
    ("Multiple ground truths: a single query can have several correct targets — every "
     "database image showing the same object under the requested modification", {"bold": True}),
    ("Curated hard negatives: right object / wrong modification, or right modification / "
     "wrong object — one modality alone is not enough", {}),
], top=Inches(5.45), size=13.5, gap=4)

# =========================================================
# 5 — i-CIR scale (restored)
# =========================================================
s = new_slide("The i-CIR benchmark — scale", notes=(
    "[0:35]  A word on scale. i-CIR contains a database of 752,092 images, 202 object "
    "instances, and 1,883 composed queries, with a median of 6 queries per instance.\n\n"
    "The histograms show the distributions: image and text queries per instance, hard "
    "negatives per instance, and the number of correct targets per composed query — "
    "confirming that many queries have more than one ground truth."))
stats = [("752,092", "database images"), ("202", "object instances"),
         ("1,883", "composed queries"), ("6", "median queries / instance")]
for i, (num, lab) in enumerate(stats):
    c = card(s, Inches(0.55 + i * 2.28), Inches(1.6), Inches(2.08), Inches(0.95),
             text=num, size=21, bold=True, color=RED)
    p2 = c.text_frame.add_paragraph()
    p2.text = lab
    p2.alignment = PP_ALIGN.CENTER
    style_runs(p2, 11.5, color=MUTED)
add_pic(s, IMG["paper_stats"], Inches(0.55), Inches(2.95), Inches(8.9), Inches(2.6))
txt(s, Inches(0.55), Inches(5.75), Inches(8.9), Inches(0.4),
    "distributions recomputed from the benchmark (cf. Psomas et al., 2025, Fig. 2)",
    size=11, color=MUTED, align=PP_ALIGN.CENTER)

# =========================================================
# 6 — zero-shot setting
# =========================================================
s = new_slide("A zero-shot setting", notes=(
    "[0:40]  In this thesis we adopt a ZERO-SHOT setting: we do not train or fine-tune "
    "any model on the i-CIR dataset.\n\n"
    "Why does this matter? Supervised composed image retrieval requires labelled triplets "
    "— a query image, a query text, and a correct target image. At the instance level, "
    "collecting these annotations for every new object would be expensive and hard to "
    "scale. Worse, a model trained on a fixed set of instances does not necessarily "
    "generalize to instances it has never seen — and in a real system new objects appear "
    "constantly.\n\n"
    "So a supervised approach would decrease both the scalability and the generalizability "
    "of the system. With frozen pre-trained models, the system applies directly to "
    "previously unseen object instances, with no additional training."))
bullets(s, [
    ("Zero-shot: we do not train or fine-tune any model on i-CIR", {"bold": True}),
    ("Supervised CIR needs labelled triplets (query image, query text, target image) — "
     "at the instance level, every new object would require new annotation", {}),
    ("That decreases scalability, and it limits generalization to instances never seen at "
     "training time — yet in practice new objects appear constantly", {}),
    ("Frozen pre-trained models apply directly to previously unseen instances, no "
     "per-instance training", {"bold": True}),
], top=Inches(1.95), size=17, gap=12)
takeaway(s, "No training anywhere — every model in this work stays frozen.", Inches(6.0))

# =========================================================
# 7 — problem formulation + evaluation
# =========================================================
s = new_slide("Problem formulation and evaluation", notes=(
    "[0:55]  Formally: given a query image and a query text, retrieve the database images "
    "relevant to both.\n\n"
    "The system assigns a relevance score to EVERY database image and sorts them from "
    "most to least relevant. This ranked list is the output; evaluation checks WHERE the "
    "correct targets land in it.\n\n"
    "For a single query we use Average Precision. AP is high when the correct images "
    "appear near the top of the ranking; if they appear lower, the score decreases.\n\n"
    "We then compute mean Average Precision by averaging AP across all queries.\n\n"
    "However, the number of queries is not the same for every instance. If we averaged "
    "directly over all queries, instances with more queries would dominate the result.\n\n"
    "For this reason our main metric is macro-mAP: we first average within each object "
    "instance, then across instances — so every instance contributes equally."))
add_pic(s, IMG["eq_problem"], Inches(0.9), Inches(1.3), Inches(8.2), Inches(0.55))
txt(s, Inches(0.75), Inches(2.0), Inches(8.6), Inches(0.32),
    "The system scores every database image and ranks them, most to least relevant:",
    size=13, color=INK)
pos = {0, 2, 5}
for i in range(8):
    card(s, Inches(0.75 + i * 0.72), Inches(2.4), Inches(0.6), Inches(0.6),
         text="✓" if i in pos else "",
         fill=RGBColor(0x0C, 0xA3, 0x0C) if i in pos else LIGHT,
         line=None, size=18, bold=True, color=WHITE)
txt(s, Inches(6.7), Inches(2.47), Inches(2.7), Inches(0.5),
    "ranked list (green = correct)", size=12, color=MUTED)
mrows = [
    ("eq_pk", "Precision@k", "fraction of the top-k results that are correct", MUTED),
    ("eq_ap", "Average Precision", "high when correct images appear near the top", MUTED),
    ("eq_map", "mAP", "AP averaged over all queries", MUTED),
    ("eq_mmap", "macro-mAP",
     "main metric — average within each instance, then across instances", RED),
]
y = Inches(3.18)
for eq, name, desc, dc in mrows:
    b = txt(s, Inches(0.75), y + Inches(0.02), Inches(3.15), Inches(0.85),
            name, size=14, bold=True, color=INK)
    p2 = b.text_frame.add_paragraph()
    p2.text = desc
    style_runs(p2, 11, color=dc)
    add_pic(s, IMG[eq], Inches(4.05), y, Inches(5.3), Inches(0.78))
    y += Inches(0.9)

# =========================================================
# 7 — vision-language models
# =========================================================
s = new_slide("Vision–language models: CLIP and SigLIP", notes=(
    "[1:10]  The main technical challenge is combining the information from the "
    "query image and the query text. To represent both modalities we use "
    "frozen, pre-trained vision–language models: CLIP and SigLIP.\n\n"
    "These use a dual-encoder (two-tower) architecture. One encoder processes images, "
    "the other text, and both map their inputs into the SAME embedding space. Because "
    "the representations share this space, we measure semantic similarity with cosine "
    "similarity: an image scores high when its embedding is close to the embedding of "
    "the corresponding text.\n\n"
    "The models are used completely frozen — we only extract embeddings, we never update "
    "or fine-tune their parameters.\n\n"
    "CLIP and SigLIP follow the same dual-encoder principle but differ in TRAINING. CLIP "
    "uses a batch-level contrastive objective based on a softmax: for each image, the "
    "matching text must be distinguished from the other texts in the batch, and vice "
    "versa. SigLIP instead uses a pairwise sigmoid loss: each image–text pair is treated "
    "independently as matching or non-matching. This avoids the global softmax and "
    "scales better to large batches.\n\n"
    "We compare these frozen backbones to test whether our method depends on one "
    "particular vision–language model."))
add_pic(s, IMG["clip"], Inches(0.5), Inches(1.5), Inches(5.3), Inches(2.7))
add_pic(s, IMG["clip_siglip"], Inches(5.95), Inches(1.5), Inches(3.5), Inches(2.7))
bullets(s, [
    ("Dual-encoder (two-tower): one image encoder, one text encoder, both mapping into "
     "the same embedding space → similarity = cosine similarity", {}),
    ("Used completely frozen: we only extract embeddings, never update parameters", {"bold": True}),
    ("CLIP — compares all image–text pairs within a batch, softmax contrastive objective", {}),
    ("SigLIP — classifies each pair independently, sigmoid loss: no global softmax, "
     "scales better to large batches", {}),
], top=Inches(4.45), size=13.5, gap=4)

# =========================================================
# 8 — BASIC
# =========================================================
s = new_slide("The baseline: BASIC", notes=(
    "[1:10]  The baseline for our work is BASIC, the zero-shot method introduced "
    "together with the i-CIR benchmark by Psomas and colleagues in 2025.\n\n"
    "BASIC processes the two parts of the composed query SEPARATELY. The query image "
    "is encoded by the visual encoder, the query text by the text encoder. The "
    "database images are encoded in advance by the visual encoder.\n\n"
    "For every database image, BASIC computes two similarity scores: the VISUAL "
    "similarity, how similar the database image is to the query image; and the "
    "TEXTUAL similarity, how well it matches the query text. Both are dot products "
    "between the corresponding embeddings.\n\n"
    "The two scores are then MULTIPLIED. This multiplicative fusion acts like a soft "
    "logical AND: a database image receives a high final score only if it matches both "
    "the query image and the query text.\n\n"
    "BASIC uses CLIP ViT-L/14 as its main backbone — the backbone most commonly used in "
    "the zero-shot CIR literature. We adopt the same backbone as our main one, so our "
    "results are directly comparable to BASIC and to the literature.\n\n"
    "The method also includes additional processing steps — centering, score "
    "normalization, contextualization — which I will discuss later."))
txt(s, Inches(0.68), Inches(0.86), Inches(8.6), Inches(0.3),
    "introduced with the i-CIR benchmark  ·  Psomas et al., 2025",
    size=12.5, italic=True, color=WHITE)
add_pic(s, IMG["basic"], Inches(0.55), Inches(1.45), Inches(8.9), Inches(3.05))
bullets(s, [
    ("The two parts of the query are processed separately, each by its own encoder", {}),
    ("For every database image: visual similarity sᵛ to the query image, and textual "
     "similarity sᵗ to the query text (dot products of embeddings)", {}),
    ("The two scores are multiplied — a soft logical AND: a high final score requires "
     "matching BOTH the query image and the query text", {"bold": True}),
    ("Main backbone: CLIP ViT-L/14 — the standard in zero-shot CIR; we adopt it too, so "
     "our results are directly comparable", {}),
], top=Inches(4.6), size=13.5, gap=4)

# =========================================================
# 9 — limitation + our idea
# =========================================================
s = new_slide("Limitation of the baseline — our proposal", notes=(
    "[1:10]  The main limitation of the baseline is that the query image and the "
    "query text are NEVER PROCESSED JOINTLY. The image branch only sees the "
    "query image, the text branch only sees the query text. As a result the model "
    "does not explicitly represent the complete composed query.\n\n"
    "To address this we introduce a THIRD scoring branch based on a generated target "
    "caption.\n\n"
    "This is the full pipeline. The query image and query text go, as before, to the "
    "image and text branches. In addition, together, they go to a multimodal language "
    "model that imagines the target and writes a caption of it. That caption is embedded "
    "by the same frozen text encoder and becomes a third similarity score.\n\n"
    "So for every database image we now compute three scores — to the image, to the text, "
    "and to the generated caption. They are normalized and multiplied, a soft logical AND "
    "across the three branches, as the figure shows. [Fusion formulas: backup slide 23.] "
    "Next, the captioner."))
txt(s, Inches(0.68), Inches(1.35), Inches(8.6), Inches(0.4),
    "In BASIC, the query image and the query text are never processed jointly — "
    "each branch sees only half of the query.",
    size=15.5, bold=True, color=RED)
txt(s, Inches(0.68), Inches(1.8), Inches(8.6), Inches(0.4),
    "Our proposal: add a THIRD branch — a caption of the imagined target, generated from "
    "the image and text together.", size=15.5, bold=True, color=INK)
add_pic(s, IMG["pipeline"], Inches(0.3), Inches(2.4), Inches(9.4), Inches(4.3))

# =========================================================
# 10 — the MLLM
# =========================================================
s = new_slide("The captioner: InternVL 3.5-8B", notes=(
    "[0:50]  For target-caption generation we use the frozen 8-billion-parameter version "
    "of InternVL 3.5 — an open-source multimodal LLM that can jointly process an image "
    "and text, which is exactly what the caption branch requires. At 8B parameters it is "
    "a practical balance of quality and cost, and we use it entirely without fine-tuning.\n\n"
    "Concretely, it receives the query image, the query text, and our instruction to "
    "describe the imagined target. In the chair query it sees the specific chair together "
    "with the text 'placed around a table', and generates: 'White metallic chair with "
    "scrolled armrests placed around a table.'\n\n"
    "Notice that this single caption contains information from BOTH parts of the query — "
    "the visual identity of the chair and the context specified by the text. That is the "
    "signal the two separate branches cannot represent on their own. We embed it with the "
    "same frozen text encoder to get the third similarity score."))
bullets(s, [
    ("Open-source multimodal LLM that jointly processes an image and a text — exactly "
     "what the caption branch requires", {}),
    ("8B parameters: a practical balance of quality and cost, competitive across "
     "vision–language benchmarks — and used entirely without fine-tuning", {"bold": True}),
], top=Inches(1.45), size=15.5, gap=7)
MY, MH = Inches(3.1), Inches(1.5)
add_pic(s, IMG["ref"], Inches(0.68), MY, Inches(1.5), MH, border=MUTED)
card(s, Inches(2.28), MY + Inches(0.22), Inches(1.55), Inches(1.05),
     text="“placed around\na table”", size=11, italic=True, color=BLUE)
arrow(s, Inches(3.93), MY + Inches(0.52), Inches(0.45))
card(s, Inches(4.48), MY + Inches(0.15), Inches(1.7), Inches(1.2), fill=ROSE,
     text="InternVL\n3.5-8B", size=12.5, bold=True)
arrow(s, Inches(6.28), MY + Inches(0.52), Inches(0.42))
card(s, Inches(6.78), MY, Inches(2.5), MH,
     text="“White metallic chair with scrolled armrests placed around a table.”",
     size=11, italic=True, align=PP_ALIGN.LEFT)
txt(s, Inches(0.68), MY + MH + Inches(0.02), Inches(3.2), Inches(0.35),
    "query image + query text, together", size=11, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, Inches(6.78), MY + MH + Inches(0.02), Inches(2.5), Inches(0.35),
    "generated target caption", size=11, color=RED, align=PP_ALIGN.CENTER)
takeaway(s, "One caption carries both the visual identity and the requested context.",
         Inches(6.15))

# =========================================================
# 11 — prompt design
# =========================================================
s = new_slide("Designing the instruction", notes=(
    "[0:45]  An important part of the caption-generation stage is the INSTRUCTION given "
    "to the multimodal model.\n\n"
    "The quality of the generated caption depends strongly on how the task is formulated. "
    "For this reason I carried out an extensive prompt-design process and tested more "
    "than twenty different instructions.\n\n"
    "I varied several elements: the level of detail requested, the use of examples, the "
    "order of the constraints, and restrictions intended to prevent the model from "
    "introducing irrelevant information.\n\n"
    "Based on these experiments I selected five complementary prompting strategies. Each "
    "encourages the model to describe the imagined target from a slightly different "
    "perspective, and together they produce a diverse set of candidate captions."))
bullets(s, [
    ("The caption quality depends strongly on how the task is formulated", {}),
    ("Extensive prompt design: more than 20 instructions tested, varying the level of "
     "detail, the use of examples, the order of constraints, and restrictions against "
     "irrelevant information", {}),
    ("Five complementary strategies selected — each describes the imagined target from a "
     "different perspective, together producing a diverse set of captions", {"bold": True}),
], top=Inches(1.5), size=16, gap=8)
strategies = ["adaptive\nprecision", "constraint\npriority", "anti-noise\ncontract",
              "slot-fill\ntemplate", "draft–refine"]
for i, name in enumerate(strategies):
    card(s, Inches(0.7 + i * 1.76), Inches(4.15), Inches(1.6), Inches(0.95),
         text=name, size=11.5, bold=True, color=INK, fill=ROSE, line=None)
txt(s, Inches(0.7), Inches(5.25), Inches(8.6), Inches(0.4),
    "the five prompting strategies — their caption embeddings are averaged (avg-5)",
    size=12, color=MUTED, align=PP_ALIGN.CENTER)
takeaway(s, "Prompt design is part of the method, not an afterthought.", Inches(5.8))

# =========================================================
# 12 — effect of the caption branch
# =========================================================
s = new_slide("Effect of the caption branch", notes=(
    "[0:55]  Adding a single generated caption improves macro-mAP by about 5.9 points, "
    "from 17.95 to 23.83. This suggests the caption captures useful compositional "
    "information not fully represented by the image and text branches separately.\n\n"
    "Using MULTIPLE captions improves the result further. We generate five captions with "
    "different prompting strategies, average their embeddings, and use the resulting "
    "representation for retrieval. This reduces the influence of prompt-specific wording "
    "or irrelevant details, and emphasizes information that is consistent across the "
    "captions.\n\n"
    "With this five-caption average, macro-mAP reaches 25.41 — about 7.5 points above the "
    "original image–text baseline.\n\n"
    "From here on, all results use CLIP ViT-L/14 as the backbone unless stated otherwise."))
table(s, [
    ("", "img × txt", "img × txt × caption\n(single)", "img × txt × caption\n(avg-5)"),
    ("macro-mAP", "17.95", "23.83", "25.41"),
    ("gain", "—", "+5.9", "+7.5"),
], Inches(0.9), Inches(2.0), Inches(8.2), row_h=0.72, size=15, highlight_rows={1})
txt(s, Inches(0.9), Inches(4.3), Inches(8.2), Inches(0.4),
    "CLIP ViT-L/14 · raw similarity product · no post-processing",
    size=12, color=MUTED, align=PP_ALIGN.CENTER)
bullets(s, [
    ("A single caption already adds +5.9 macro-mAP — compositional information the two "
     "branches miss on their own", {}),
    ("Averaging five captions removes prompt-specific wording and keeps what is "
     "consistent across them: +7.5 in total", {"bold": True}),
], top=Inches(4.85), size=15, gap=6)
txt(s, Inches(0.68), Inches(6.35), Inches(8.6), Inches(0.35),
    "From here on: CLIP ViT-L/14 as the vision–language backbone, unless stated otherwise.",
    size=12, italic=True, color=RED, align=PP_ALIGN.CENTER)

# =========================================================
# 14 — post-processing (centering + contextualisation)
# =========================================================
s = new_slide("Two refinements: centering and contextualisation", notes=(
    "[0:50]  Two further techniques improve retrieval — one on the embeddings, one on the "
    "query text itself.\n\n"
    "CENTERING, applied to the embeddings. Vision–language embedding spaces contain a "
    "large common component shared by many embeddings — generic visual or linguistic "
    "patterns rather than the distinctive information needed for retrieval. We compute a "
    "mean embedding for each modality on a separate corpus, subtract it, and re-normalise. "
    "This isolates the semantic content from the shared distributional bias.\n\n"
    "CONTEXTUALISATION, applied to the raw query text — so it is a query-side step, not a "
    "score post-processing step. The text encoder is trained mostly on full captions, so a "
    "bare fragment such as 'sculpture' or 'during sunset' is out-of-distribution and "
    "produces a text feature poorly aligned with image features. We therefore enrich the "
    "query with terms from a subject corpus, before and after — 'dog during the sunset', "
    "'sculpture dog' — then embed and average the variants, giving a more robust textual "
    "feature."))
txt(s, Inches(0.68), Inches(1.45), Inches(4.3), Inches(0.4), "Centering",
    size=16, bold=True, color=RED)
txt(s, Inches(0.68), Inches(1.85), Inches(4.3), Inches(0.32),
    "on the embeddings", size=11.5, italic=True, color=MUTED)
add_pic(s, IMG["eq_centering"], Inches(0.68), Inches(2.25), Inches(4.3), Inches(0.75))
bullets(s, [
    ("Embeddings share a large common component: generic visual / linguistic patterns, "
     "not distinctive information", {}),
    ("Subtract a per-modality mean (computed on a separate corpus) and re-normalise", {"bold": True}),
], left=Inches(0.68), top=Inches(3.15), w=Inches(4.3), size=12.5, gap=5)
txt(s, Inches(5.15), Inches(1.45), Inches(4.3), Inches(0.4), "Contextualisation",
    size=16, bold=True, color=RED)
txt(s, Inches(5.15), Inches(1.85), Inches(4.3), Inches(0.32),
    "on the raw query text (a query-side step)", size=11.5, italic=True, color=MUTED)
bullets(s, [
    ("The text encoder is trained on full captions; a bare fragment is "
     "out-of-distribution", {}),
    ("Enrich the query text with terms from a subject corpus, before and after — then "
     "embed and average the variants", {"bold": True}),
], left=Inches(5.15), top=Inches(2.25), w=Inches(4.3), size=12.5, gap=5)
for i, (raw, ctx) in enumerate([("“during sunset”", "“dog during the sunset”"),
                                ("“sculpture”", "“sculpture dog”")]):
    card(s, Inches(5.15), Inches(3.65 + i * 0.8), Inches(1.75), Inches(0.62),
         text=raw, size=10.5, italic=True, color=BLUE)
    txt(s, Inches(6.95), Inches(3.72 + i * 0.8), Inches(0.35), Inches(0.45), "→",
        size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)
    card(s, Inches(7.35), Inches(3.65 + i * 0.8), Inches(1.95), Inches(0.62),
         text=ctx, size=10.5, italic=True, color=INK)
takeaway(s, "Centering — remove what is common to everything.\n"
            "Contextualisation — speak to the encoder in its own language.",
         Inches(5.55))

# =========================================================
# 15 — the ladder
# =========================================================
s = new_slide("Adding the steps on top of each other", notes=(
    "[0:40]  Here we add each step cumulatively, for the baseline in blue "
    "and our caption pipeline in red, on CLIP ViT-L/14.\n\n"
    "The caption pipeline stays above the baseline along the entire ladder. Centering "
    "gives a large jump for the baseline; contextualisation gives the largest single "
    "gain.\n\n"
    "[If asked: two further BASIC steps — the contrastive projection and query expansion — "
    "do not help once the caption branch is present, so they are omitted. Numbers on the "
    "full-ablation backup slide.]"))
add_pic(s, IMG["ladder_l"], Inches(0.55), Inches(1.5), Inches(8.9), Inches(4.5))
takeaway(s, "The caption pipeline stays ahead at every step of the ladder.", Inches(6.2))

# =========================================================
# 16 — backbones (table + results)
# =========================================================
s = new_slide("Does it transfer? Four frozen backbones", notes=(
    "[0:55]  We evaluate the pipeline with four frozen backbones spanning two model "
    "families and two capacity levels — CLIP ViT-L/14 (the backbone most commonly reported "
    "in the zero-shot CIR literature, which makes our numbers comparable), the "
    "higher-capacity CLIP ViT-H/14, SigLIP ViT-L/16, and the more recent SigLIP2 g/16, "
    "which is the best-aligned encoder of the four.\n\n"
    "The caption branch improves EVERY backbone. The largest jump is on SigLIP-L, where "
    "macro-mAP roughly doubles, from 21.1 to 43.3. SigLIP2 is the strongest base model.\n\n"
    "This shows the contribution does not depend on one particular vision–language model."))
add_pic(s, IMG["backbones"], Inches(0.55), Inches(1.45), Inches(5.55), Inches(4.4))
table(s, [
    ("backbone", "dim"),
    ("CLIP ViT-L/14", "768"),
    ("CLIP ViT-H/14", "1024"),
    ("SigLIP ViT-L/16", "1024"),
    ("SigLIP2 g/16", "1536"),
], Inches(6.3), Inches(2.0), Inches(3.1), row_h=0.5, size=12, highlight_rows={4})
txt(s, Inches(6.3), Inches(4.6), Inches(3.1), Inches(1.2),
    "Two families — softmax-contrastive (CLIP) and sigmoid (SigLIP) — at two capacity "
    "levels.", size=12, color=MUTED)
takeaway(s, "The caption branch helps every backbone — it is not a CLIP-specific trick.",
         Inches(6.2))

# =========================================================
# 17 — best overall system
# =========================================================
s = new_slide("Our best overall system", notes=(
    "[0:50]  To state the final system plainly: the SigLIP2 backbone, the caption branch "
    "with five averaged captions, and post-processing through centering and "
    "contextualisation. Every component frozen.\n\n"
    "61.98 macro-mAP and 62.07 mAP. Beyond mAP: the correct target is ranked first for 71% "
    "of queries, and appears in the top ten for 95% of queries. Recall@100 is 94%.\n\n"
    "That high recall is important — I come back to it in the future work."))
recipe = [("backbone", "SigLIP2 g/16-384"), ("caption branch", "InternVL 3.5-8B\navg-5"),
          ("post-processing", "centering\n+ contextualisation")]
for i, (lab, val) in enumerate(recipe):
    c = card(s, Inches(0.75 + i * 2.95), Inches(1.6), Inches(2.7), Inches(1.3),
             text=lab, size=11.5, color=MUTED)
    for ln in val.split("\n"):
        p2 = c.text_frame.add_paragraph()
        p2.text = ln
        p2.alignment = PP_ALIGN.CENTER
        style_runs(p2, 13.5, bold=True, color=INK)
    if i < 2:
        txt(s, Inches(3.42 + i * 2.95), Inches(2.0), Inches(0.4), Inches(0.5),
            "+", size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)
metrics = [("61.98", "macro-mAP"), ("62.07", "mAP"), ("71.1%", "Hit@1"),
           ("94.9%", "Hit@10"), ("93.8%", "Recall@100")]
for i, (num, lab) in enumerate(metrics):
    hero = (i == 0)
    c = card(s, Inches(0.62 + i * 1.79), Inches(3.35), Inches(1.62), Inches(1.15),
             text=num, size=20 if hero else 17, bold=True,
             color=RED if hero else INK, fill=ROSE if hero else LIGHT, line=None)
    p2 = c.text_frame.add_paragraph()
    p2.text = lab
    p2.alignment = PP_ALIGN.CENTER
    style_runs(p2, 11, color=MUTED)
bullets(s, [
    ("+4.3 macro-mAP over the strongest BASIC configuration (57.69)", {}),
    ("The correct target reaches the top-10 for 94.9% of queries — most of the remaining "
     "error is in the ORDER of already-retrieved candidates", {"bold": True}),
], top=Inches(4.75), size=15, gap=6)

# =========================================================
# 19 — qualitative success
# =========================================================
s = new_slide("Qualitative example — success", notes=(
    "[0:30, flex]  Homer Simpson, 'as a plushie'. The plain image × text product scores "
    "only 4% AP — it drifts to generic plushies: SpongeBob, minions, a teddy bear. BASIC "
    "reaches 33%. Adding the caption — 'Homer Simpson as a plushie with soft texture and "
    "the same facial features' — pins the retrieval to actual Homer plushies and lifts AP "
    "to 88%. This is the caption branch doing exactly what it is meant to do: keeping the "
    "specific identity while honouring the modification."))
add_pic(s, IMG["qual_success"], Inches(0.4), Inches(1.5), Inches(9.2), Inches(5.0))

# =========================================================
# 20 — qualitative failure
# =========================================================
s = new_slide("Qualitative example — failure", notes=(
    "[0:30, flex]  Gold metal eyeglasses, 'worn by a dog'. Here all three methods fail — "
    "AP near zero. The single ground truth is a photo of a dog actually wearing these "
    "glasses. Every method instead returns glasses photographed on tables or held up, "
    "because the visual branch anchors hard on the glasses themselves and the caption "
    "cannot overcome a target this rare.\n\n"
    "This is the characteristic failure mode: when the modification demands a drastic, "
    "rarely-photographed context, the image branch dominates and the correct target is "
    "simply not near the top."))
add_pic(s, IMG["qual_failure"], Inches(0.4), Inches(1.5), Inches(9.2), Inches(5.0))

# =========================================================
# 21 — limitations & future work
# =========================================================
s = new_slide("Limitations and future work", notes=(
    "[0:50]  Limitations. The main cost is caption generation — about 15 seconds per "
    "query, which dominates online inference. The results also depend on prompt design: "
    "averaging five captions mitigates this but does not remove it. And the fixed, "
    "unweighted product cannot adapt when one branch should matter more.\n\n"
    "Future work. A stronger or more efficient captioner — or a single high-quality caption "
    "that matches avg-5 without five generations — would cut the main compute cost. The "
    "fixed product could be replaced by a learned or weighted fusion of the three "
    "similarities, relaxing the zero-shot constraint in exchange for accuracy.\n\n"
    "Because recall is already high — the correct instance lands in the top-ranked "
    "shortlist for the large majority of queries — a second-stage RERANKER that reorders "
    "this shortlist is a promising direction: most of the remaining error is in the "
    "ordering of candidates that are already retrieved. A text-conditioned reranker that "
    "reorders by similarity to the target caption, rather than to the query image, is "
    "a natural candidate, since it keeps the modification in the loop at the reranking "
    "stage.\n\n"
    "Finally, the caption-as-third-modality idea is benchmark-agnostic and could be "
    "evaluated beyond i-CIR, on category-level CIR benchmarks and other instance-level "
    "retrieval settings."))
txt(s, Inches(0.68), Inches(1.45), Inches(8.6), Inches(0.4), "Limitations",
    size=16, bold=True, color=RED)
bullets(s, [
    ("Compute — caption generation dominates online inference (≈ 15 s per query)", {}),
    ("Prompt sensitivity — results depend on instruction design; avg-5 mitigates but does "
     "not remove it", {}),
    ("The fixed, unweighted product cannot re-weight the branches per query", {}),
], top=Inches(1.9), size=14.5, gap=4)
txt(s, Inches(0.68), Inches(3.7), Inches(8.6), Inches(0.4), "Future work",
    size=16, bold=True, color=RED)
bullets(s, [
    ("A stronger or more efficient captioner — or a single caption matching avg-5 without "
     "five generations — would cut the main compute cost", {}),
    ("A learned or weighted fusion of the three similarities, relaxing the zero-shot "
     "constraint in exchange for accuracy", {}),
    ("Second-stage reranking: recall is already high (Hit@10 = 94.9%), so most of the "
     "remaining error is in the ORDER of candidates already retrieved — a text-conditioned "
     "reranker keeps the modification in the loop", {"bold": True}),
    ("The caption-as-third-modality idea is benchmark-agnostic: category-level CIR and "
     "other instance-level retrieval settings", {}),
], top=Inches(4.15), size=14.5, gap=4)

# =========================================================
# 22 — conclusions
# =========================================================
s = new_slide("Conclusions", notes=(
    "[0:40]  Two conclusions.\n\n"
    "First, the caption is the single most effective component. Adding the target caption "
    "improves every backbone on the full dataset.\n\n"
    "Second, SigLIP2 is the strongest backbone for the pipeline, reaching the thesis-best "
    "61.98 macro-mAP.\n\n"
    "[If asked to go further: the caption alone already beats the plain image × text "
    "product on CLIP-L, and centering and the caption are complementary — backup slides "
    "have the numbers.]"))
bullets(s, [
    ("The caption is the single most effective component.", {"bold": True, "size": 20}),
    ("Adding the target caption improves every backbone on the full dataset.", {"size": 17}),
    ("", {"nobullet": True, "size": 8}),
    ("SigLIP2 is the strongest backbone for the pipeline.", {"bold": True, "size": 20}),
    ("61.98 macro-mAP — the best result in the thesis, fully zero-shot.", {"size": 17}),
], top=Inches(2.1), size=17, gap=14)

# =========================================================
# 23 — thank you (full-bleed red)
# =========================================================
s = new_slide("")
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = RED
bg.line.fill.background()
tf = bg.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Thank you!"
p.alignment = PP_ALIGN.CENTER
style_runs(p, 54, bold=True, color=WHITE)
s.notes_slide.notes_text_frame.text = "Thank the committee; open for questions."

# =========================================================
# backups
# =========================================================
s = new_slide("Backup — combining the three scores", notes=(
    "How the three similarity scores are fused. Min-based normalization puts the branches "
    "on comparable scales; negative residuals are clamped to zero — no positive evidence. "
    "The normalized scores are multiplied, a soft logical AND: a good result must match "
    "all three branches at once.\n\n"
    "On top, a Harris-inspired penalty — from the Harris corner detector, where a strong "
    "response in one direction only is not enough. A candidate must not win on one branch "
    "alone: the exact chair not around a table, or chairs around a table that are not the "
    "query chair, are both suppressed. The penalty is weighted by a hyperparameter lambda, "
    "fixed across all experiments.\n\n"
    "The surface shows the effect: the score is high only when both arguments are "
    "simultaneously large."))
txt(s, Inches(0.68), Inches(1.45), Inches(4.7), Inches(0.35),
    "1. Min-normalise, then clamp negatives to zero", size=14, bold=True, color=RED)
add_pic(s, IMG["eq_minnorm"], Inches(0.7), Inches(1.9), Inches(4.4), Inches(0.7))
bullets(s, [
    ("Puts the three branches on comparable scales; a negative residual = no positive "
     "evidence for that candidate", {}),
], left=Inches(0.68), top=Inches(2.7), w=Inches(4.7), size=13, gap=4)
txt(s, Inches(0.68), Inches(3.75), Inches(4.7), Inches(0.35),
    "2. Multiply, with a Harris-inspired penalty (λ)", size=14, bold=True, color=RED)
add_pic(s, IMG["eq_harris"], Inches(0.7), Inches(4.2), Inches(4.4), Inches(0.8))
bullets(s, [
    ("A candidate must be strong on ALL branches — the exact chair NOT around a table, or "
     "chairs around a table that are NOT the query chair, are both suppressed", {"bold": True}),
], left=Inches(0.68), top=Inches(5.15), w=Inches(4.7), size=13, gap=4)
add_pic(s, IMG["harris"], Inches(5.5), Inches(1.55), Inches(4.0), Inches(4.4))
txt(s, Inches(5.5), Inches(5.95), Inches(4.0), Inches(0.35),
    "high only when both scores are large", size=11.5, italic=True, color=MUTED,
    align=PP_ALIGN.CENTER)

s = new_slide("Backup — full CLIP-L ablation", notes=(
    "Every rung of the ladder, for the three pipelines. macro-mAP, full i-CIR, CLIP "
    "ViT-L/14. Projection and query expansion are shown here: they do not help once the "
    "caption branch is present, which is why the pipeline stops at contextualisation."))
table(s, [
    ("step", "img × txt", "× caption (single)", "× caption (avg-5)"),
    ("raw product", "17.95", "23.83", "25.41"),
    ("+ centering", "27.76", "31.02", "31.87"),
    ("+ sim-normalisation", "27.40", "30.28", "31.97"),
    ("+ Harris criterion", "28.33", "31.25", "32.82"),
    ("+ contextualisation", "32.25", "34.06", "35.76"),
    ("+ projection", "32.48", "32.58", "33.80"),
    ("+ query expansion", "30.28", "31.53", "32.39"),
], Inches(1.0), Inches(1.75), Inches(8.0), row_h=0.5, size=13.5, highlight_rows={5})
txt(s, Inches(1.0), Inches(5.5), Inches(8.0), Inches(0.4),
    "macro-mAP · full i-CIR · CLIP ViT-L/14 · the pipeline stops at contextualisation",
    size=12, color=MUTED, align=PP_ALIGN.CENTER)

s = new_slide("Backup — SigLIP2 ladder", notes=(
    "The best backbone, with the full ladder — the slide we skipped in the talk. The shape "
    "is the same as on CLIP-L: the caption pipeline stays above BASIC throughout, and "
    "contextualisation gives the final jump. Best configuration 61.98 macro-mAP, against "
    "57.69 for BASIC at the same rung, and 38.1 for the plain image × text product."))
add_pic(s, IMG["ladder_s2"], Inches(0.55), Inches(1.5), Inches(8.9), Inches(4.6))
txt(s, Inches(0.55), Inches(6.2), Inches(8.9), Inches(0.4),
    "macro-mAP · full i-CIR · SigLIP2 g/16-384", size=12, color=MUTED,
    align=PP_ALIGN.CENTER)

s = new_slide("Backup — component contributions (CLIP-L)", notes=(
    "Each branch on its own, then the product, then the triplet. Raw similarity, no "
    "post-processing, macro-mAP, CLIP ViT-L/14. Two things to note: the caption alone "
    "(19.51) already beats the full image × text product (17.95); and adding the caption "
    "to the product lifts it from 17.95 to 25.41."))
table(s, [
    ("configuration", "macro-mAP"),
    ("image only", "2.75"),
    ("text only", "3.26"),
    ("caption only (avg-5)", "19.51"),
    ("image × text", "17.95"),
    ("image × text × caption (avg-5)", "25.41"),
], Inches(2.15), Inches(1.8), Inches(5.7), row_h=0.55, size=15, highlight_rows={5})
bullets(s, [
    ("The caption alone (19.51) already beats the image × text product (17.95)", {"bold": True}),
    ("Adding the caption to the product: 17.95 → 25.41 (+7.5)", {}),
], top=Inches(5.5), left=Inches(2.15), w=Inches(5.7), size=13, gap=5)
txt(s, Inches(2.15), Inches(6.42), Inches(5.7), Inches(0.3),
    "raw similarity · no post-processing · macro-mAP · CLIP ViT-L/14",
    size=11.5, color=MUTED, align=PP_ALIGN.CENTER)

s = new_slide("Backup — centering and the caption are complementary", notes=(
    "The caption gain is not absorbed by centering: it survives on top of it. At the raw "
    "product the caption adds +7.5; after centering it still adds +4.1. The two "
    "contributions stack, which is why the best configuration keeps both. CLIP-L, macro-mAP."))
table(s, [
    ("", "img × txt", "× caption (avg-5)", "Δ caption"),
    ("raw product", "17.95", "25.41", "+7.5"),
    ("+ centering", "27.76", "31.87", "+4.1"),
], Inches(1.3), Inches(2.0), Inches(7.4), row_h=0.62, size=15, highlight_rows={2})
bullets(s, [
    ("The caption gain survives centering — the two contributions stack rather than "
     "overlap", {"bold": True}),
    ("So the best configuration combines both", {}),
], top=Inches(4.3), size=15, gap=7)

s = new_slide("Backup — reproducing BASIC", notes=(
    "The i-CIR paper reports 34.35 for BASIC (CLIP-L, full stack). Our re-run of the "
    "released code gives 32.48. The gap is concentrated in the image branch and is "
    "consistent with re-extracting features in a newer software stack. Both methods are "
    "evaluated in the same environment, so the comparisons are unaffected."))
table(s, [
    ("BASIC, CLIP ViT-L/14", "macro-mAP"),
    ("reported in the i-CIR paper", "34.35"),
    ("our reproduction (released code)", "32.48"),
], Inches(2.2), Inches(1.9), Inches(5.6), row_h=0.6, size=15)
bullets(s, [
    ("The gap is concentrated in the IMAGE branch — features re-extracted in a newer "
     "software stack (image decoding / resizing differ across library versions)", {}),
    ("Contextualisation also draws fresh object nouns per query — small run-to-run "
     "variance", {}),
    ("Both BASIC and our method are evaluated in the SAME environment, so all relative "
     "comparisons in this thesis are unaffected", {"bold": True}),
], top=Inches(4.1), size=15, gap=7)

s = new_slide("Backup — inference latency", notes=(
    "100 queries, RTX-6000-Ada 48 GB, five captions per query. Caption generation is "
    "≈99.7% of the online cost; retrieval itself is negligible. Database embeddings are "
    "precomputed offline."))
table(s, [
    ("stage", "CLIP-L", "SigLIP2"),
    ("caption generation (5 captions)", "≈ 14.7 s", "≈ 14.7 s"),
    ("embedding the query triplet", "22 ms", "46 ms"),
    ("total, per query (median)", "14.75 s", "14.75 s"),
], Inches(1.4), Inches(1.9), Inches(7.2), row_h=0.58, size=14, highlight_rows={3})
bullets(s, [
    ("Caption generation is ≈ 99.7% of online inference — the backbone choice does not "
     "change latency", {"bold": True}),
    ("Database embeddings are precomputed offline; ranking itself is instantaneous", {}),
    ("Levers: fewer captions, shorter generations, a smaller or quantised MLLM", {}),
], top=Inches(4.5), size=15, gap=6)

s = new_slide("Backup — the five instructions, sample captions", notes=(
    "The five prompting strategies on the tintin example, 'as an iconic rooftop sign'. "
    "Each phrases the same target differently; their embeddings are averaged (avg-5)."))
add_pic(s, IMG["tintin_q"], Inches(0.68), Inches(1.7), Inches(2.0), Inches(1.7), border=MUTED)
txt(s, Inches(0.68), Inches(3.55), Inches(2.0), Inches(0.6),
    "“as an iconic rooftop sign”", size=12, italic=True, color=BLUE,
    align=PP_ALIGN.CENTER)
caps = [("adaptive precision", "A character with a beige coat and blue shirt, as an iconic rooftop sign"),
        ("constraint-priority", "A character wearing a trench coat, blonde hair, and a blue shirt, is an iconic rooftop sign."),
        ("anti-noise contract", "Character in a beige trench coat, blue shirt, with blushing cheeks and short orange hair as an iconic rooftop sign."),
        ("slot-fill template", "a cartoon character wearing a trench coat, blushing, with short hair as an iconic rooftop sign."),
        ("draft–refine", "Cartoon character in a beige coat and blue shirt appears as an iconic rooftop sign.")]
for i, (tag, cap) in enumerate(caps):
    card(s, Inches(3.0), Inches(1.6 + i * 1.0), Inches(6.4), Inches(0.9),
         text=f"{tag}:  “{cap}”", size=10.5, align=PP_ALIGN.LEFT)

s = new_slide("Backup — the five instructions, explained", notes=(
    "What each of the five prompting strategies actually asks the model to do. They were "
    "selected from more than twenty candidates because they are complementary — each "
    "controls a different failure mode of caption generation."))
strat = [
    ("adaptive precision", "Adjusts the level of visual detail to the query — terse when "
     "the modification is simple, richer when the target is complex."),
    ("constraint-priority", "States the modification FIRST, then the identity — so the "
     "requested change is never dropped in favour of describing the object."),
    ("anti-noise contract", "Explicitly forbids inventing attributes not visible in the "
     "image — suppresses hallucinated colours, text, and background."),
    ("slot-fill template", "Fills a fixed template (object · attributes · modification) — "
     "keeps captions uniform and comparable across queries."),
    ("draft–refine", "Writes a first caption, then rewrites it once — a self-correction "
     "pass that tends to give the cleanest single sentence."),
]
for i, (tag, desc) in enumerate(strat):
    y = Inches(1.55 + i * 1.02)
    card(s, Inches(0.7), y, Inches(2.35), Inches(0.9), text=tag, size=13, bold=True,
         color=INK, fill=ROSE, line=None)
    txt(s, Inches(3.2), y + Inches(0.06), Inches(6.2), Inches(0.85), desc,
        size=12.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
